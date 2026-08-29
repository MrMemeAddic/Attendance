"""
Jagte Raho — Final Year Presentation Generator
Creates a professional 10-slide PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour Palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_CARD      = RGBColor(0x16, 0x2A, 0x3E)   # card navy
ACCENT_TEAL  = RGBColor(0x00, 0xD4, 0xB8)   # teal accent
ACCENT_GOLD  = RGBColor(0xF5, 0xA6, 0x23)   # warm gold
TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT   = RGBColor(0xC8, 0xD8, 0xE8)
TEXT_DIM     = RGBColor(0x7A, 0x96, 0xB2)
ACCENT_RED   = RGBColor(0xFF, 0x4D, 0x6D)   # alert / anti-spoof
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)   # success


# ── Helpers ───────────────────────────────────────────────────────────────────

def rgb_hex(r, g, b):
    return RGBColor(r, g, b)


def set_slide_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=TEXT_WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=16, color=TEXT_WHITE, align=PP_ALIGN.LEFT,
                       bold=False, line_spacing_pt=None, font_name="Calibri"):
    """Add a textbox with multiple lines, each as a paragraph."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def gradient_rect(slide, left, top, width, height, color1, color2):
    """Add a shape with a simple two-color gradient (left→right)."""
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 0
    gs = fill.gradient_stops
    gs[0].position = 0
    gs[0].color.rgb = color1
    gs[1].position = 1
    gs[1].color.rgb = color2
    shape.line.fill.background()
    return shape


def accent_bar(slide, left=0, top=0, width=10, height=0.04, color=ACCENT_TEAL):
    add_rect(slide, left, top, width, height, color)


# ── Slide Builders ────────────────────────────────────────────────────────────

def slide1_title(prs):
    """Slide 1 – Title / Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Gradient banner across top
    gradient_rect(slide, 0, 0, 10, 2.8, BG_CARD, RGBColor(0x0A, 0x28, 0x45))

    # Decorative teal accent bar
    accent_bar(slide, 0, 2.8, 10, 0.06, ACCENT_TEAL)

    # Hexagon emoji / icon placeholder
    add_text(slide, "⬡", 0.3, 0.25, 1.2, 1.2, font_size=52, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Project name
    add_text(slide, "JAGTE RAHO", 1.4, 0.22, 8, 0.9,
             font_size=46, bold=True, color=TEXT_WHITE, font_name="Calibri")

    # Subtitle
    add_text(slide, "Real-Time Face Recognition Attendance System",
             1.4, 1.1, 7.8, 0.6, font_size=22, color=ACCENT_TEAL, font_name="Calibri")

    # Tagline
    add_text(slide, "\"Stay Vigilant\"  —  Offline  •  Secure  •  Intelligent",
             1.4, 1.75, 7.8, 0.5, font_size=14, italic=True, color=TEXT_DIM, font_name="Calibri")

    # Divider
    accent_bar(slide, 0.5, 3.5, 9, 0.02, TEXT_DIM)

    # Tech-stack row
    techs = ["Python 3.14", "PyTorch 2.11", "MediaPipe", "FaceNet", "OpenCV", "SQLite"]
    cols = len(techs)
    col_w = 9.0 / cols
    for i, t in enumerate(techs):
        x = 0.5 + i * col_w
        add_rect(slide, x + 0.05, 3.65, col_w - 0.1, 0.42, BG_CARD)
        add_text(slide, t, x + 0.08, 3.68, col_w - 0.16, 0.36,
                 font_size=11, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Bottom info bar
    add_rect(slide, 0, 6.5, 10, 1.0, RGBColor(0x07, 0x12, 0x1E))
    add_text(slide, "Final Year Project Presentation  |  B.Tech Computer Science",
             0.3, 6.58, 9, 0.4, font_size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "Presented by: Harsh",
             0.3, 6.95, 9, 0.4, font_size=13, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def slide2_problem(prs):
    """Slide 2 – Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Top bar
    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_GOLD)
    add_text(slide, "🔍  PROBLEM STATEMENT", 0.3, 0.2, 9, 0.75,
             font_size=28, bold=True, color=TEXT_WHITE)

    problems = [
        ("📋  Manual Attendance", "Traditional roll-call is time-consuming, error-prone, and easily manipulated through proxy attendance."),
        ("🕵️  Photo Spoofing", "Simple face-detection systems are fooled by printed photos or digital screen displays — no way to detect a real person."),
        ("☁️  Privacy Concerns", "Cloud-based solutions transmit biometric data externally — unacceptable for institutions with strict data-privacy requirements."),
        ("⏱️  Scalability Issues", "Large classes (100+ students) make manual or swipe-card based systems impractically slow."),
    ]

    for i, (title, desc) in enumerate(problems):
        row = 1.35 + i * 1.3
        add_rect(slide, 0.3, row, 9.4, 1.15, BG_CARD)
        # coloured left stripe
        stripe_color = [ACCENT_GOLD, ACCENT_RED, ACCENT_TEAL, ACCENT_GREEN][i]
        add_rect(slide, 0.3, row, 0.12, 1.15, stripe_color)
        add_text(slide, title, 0.55, row + 0.05, 4, 0.42,
                 font_size=15, bold=True, color=TEXT_WHITE)
        add_text(slide, desc, 0.55, row + 0.5, 8.9, 0.6,
                 font_size=12, color=TEXT_LIGHT)


def slide3_objectives(prs):
    """Slide 3 – Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_TEAL)
    add_text(slide, "🎯  OBJECTIVES", 0.3, 0.2, 9, 0.75,
             font_size=28, bold=True, color=TEXT_WHITE)

    objectives = [
        ("01", "Automated Attendance", "Eliminate manual roll-call through real-time face recognition with high accuracy."),
        ("02", "Anti-Spoofing Security", "Implement Eye Aspect Ratio (EAR) liveness detection to prevent photo/video spoofing attacks."),
        ("03", "Offline & Private", "Store all biometric data locally in SQLite — zero cloud dependency, full data sovereignty."),
        ("04", "Audit-Ready Reports", "Auto-generate daily CSV and styled XLSX attendance reports for institutional record keeping."),
        ("05", "User-Friendly Interface", "Deliver a polished 3-tab desktop GUI for registration, recognition, and management."),
    ]

    for i, (num, title, desc) in enumerate(objectives):
        row = 1.35 + i * 1.08
        # number badge
        add_rect(slide, 0.3, row + 0.08, 0.55, 0.55, ACCENT_TEAL)
        add_text(slide, num, 0.3, row + 0.08, 0.55, 0.55,
                 font_size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.0, row + 0.05, 3.2, 0.4,
                 font_size=14, bold=True, color=ACCENT_TEAL)
        add_text(slide, desc, 1.0, row + 0.42, 8.6, 0.55,
                 font_size=11, color=TEXT_LIGHT)
        # thin divider
        if i < len(objectives) - 1:
            accent_bar(slide, 0.3, row + 0.98, 9.4, 0.01, TEXT_DIM)


def slide4_architecture(prs):
    """Slide 4 – System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_TEAL)
    add_text(slide, "🏗️  SYSTEM ARCHITECTURE", 0.3, 0.2, 9, 0.75,
             font_size=28, bold=True, color=TEXT_WHITE)

    # Pipeline flow (horizontal)
    steps = [
        ("📷", "Camera\nFrame", ACCENT_TEAL),
        ("👤", "BlazeFace\nDetection", ACCENT_GOLD),
        ("🧠", "FaceNet\nEmbedding", ACCENT_TEAL),
        ("📐", "Cosine\nSimilarity", ACCENT_GOLD),
        ("✅", "Identity +\nConfidence", ACCENT_GREEN),
    ]

    box_w = 1.55
    box_h = 1.1
    start_x = 0.25
    y = 1.55

    for i, (icon, label, color) in enumerate(steps):
        x = start_x + i * (box_w + 0.38)
        add_rect(slide, x, y, box_w, box_h, BG_CARD, line_color=color, line_width=Pt(1.5))
        add_text(slide, icon, x, y + 0.05, box_w, 0.45, font_size=22, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, y + 0.5, box_w, 0.58,
                 font_size=10, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Arrow
        if i < len(steps) - 1:
            ax = x + box_w + 0.04
            add_text(slide, "▶", ax, y + 0.3, 0.35, 0.45,
                     font_size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Module grid
    add_text(slide, "Core Modules", 0.3, 2.9, 4, 0.4,
             font_size=14, bold=True, color=ACCENT_TEAL)
    add_text(slide, "Data Layer", 5.2, 2.9, 4, 0.4,
             font_size=14, bold=True, color=ACCENT_GOLD)

    modules_l = ["detector.py  —  BlazeFace + FaceLandmarker",
                 "embedder.py  —  FaceNet 512-dim embeddings",
                 "liveness.py   —  EAR blink analysis",
                 "tracker.py    —  Multi-face frame tracking",
                 "matcher.py  —  Cosine similarity matching"]

    modules_r = ["database.py  —  SQLite CRUD operations",
                 "attendance.py  —  CSV + XLSX logging",
                 "generate_report.py  —  Summary reports",
                 "data/faces.db  —  Embedded local DB",
                 "data/attendance/  —  Daily export files"]

    add_multiline_text(slide, modules_l, 0.3, 3.3, 4.7, 2.8,
                       font_size=11, color=TEXT_LIGHT)
    add_multiline_text(slide, modules_r, 5.2, 3.3, 4.6, 2.8,
                       font_size=11, color=TEXT_LIGHT)

    # vertical divider
    add_rect(slide, 4.95, 2.85, 0.02, 3.6, TEXT_DIM)


def slide5_face_recognition(prs):
    """Slide 5 – Face Recognition Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_TEAL)
    add_text(slide, "🧠  FACE RECOGNITION PIPELINE", 0.3, 0.2, 9, 0.75,
             font_size=28, bold=True, color=TEXT_WHITE)

    # Two-stage detection box
    add_rect(slide, 0.25, 1.3, 4.55, 2.0, BG_CARD, line_color=ACCENT_TEAL, line_width=Pt(1))
    add_text(slide, "⚡  Two-Stage Detection", 0.4, 1.35, 4.2, 0.5,
             font_size=14, bold=True, color=ACCENT_TEAL)
    add_multiline_text(slide, [
        "Stage 1 — BlazeFace (fast, ~30-60 ms/frame)",
        "  → Locates face bounding boxes in real time",
        "",
        "Stage 2 — FaceLandmarker (full 478 landmarks)",
        "  → Activated only when needed for liveness",
    ], 0.4, 1.82, 4.3, 1.45, font_size=11, color=TEXT_LIGHT)

    # FaceNet box
    add_rect(slide, 5.2, 1.3, 4.55, 2.0, BG_CARD, line_color=ACCENT_GOLD, line_width=Pt(1))
    add_text(slide, "🔷  FaceNet Embeddings", 5.35, 1.35, 4.2, 0.5,
             font_size=14, bold=True, color=ACCENT_GOLD)
    add_multiline_text(slide, [
        "Model: InceptionResnetV1 (pretrained on VGGFace2)",
        "Output: 512-dimensional embedding vector",
        "",
        "Matching: Cosine Similarity",
        "  → Threshold tuned per deployment context",
        "  → Embeddings stored as BLOBs in SQLite",
    ], 5.35, 1.82, 4.3, 1.45, font_size=11, color=TEXT_LIGHT)

    # Database schema
    add_rect(slide, 0.25, 3.55, 9.5, 2.1, BG_CARD)
    add_text(slide, "🗃️  Database Schema (SQLite)", 0.4, 3.6, 5, 0.45,
             font_size=14, bold=True, color=ACCENT_TEAL)

    tables = [
        ("persons", "id, name, label, created_at"),
        ("embeddings", "id, person_id → persons, embedding BLOB, source, created_at"),
        ("recognition_log", "id, person_id, name, confidence, liveness, timestamp"),
    ]
    for i, (tbl, cols) in enumerate(tables):
        row = 4.1 + i * 0.48
        add_text(slide, tbl, 0.4, row, 2.2, 0.42,
                 font_size=12, bold=True, color=ACCENT_GOLD)
        add_text(slide, cols, 2.7, row, 6.8, 0.42,
                 font_size=11, color=TEXT_LIGHT)

    # Key metric callout
    add_rect(slide, 0.25, 5.85, 9.5, 0.8, RGBColor(0x00, 0x44, 0x3C))
    add_text(slide, "✔  Non-blocking I/O: Recognition logs written asynchronously — no UI stutter during high-traffic sessions",
             0.45, 5.92, 9.1, 0.55, font_size=12, color=ACCENT_GREEN, bold=True)


def slide6_liveness(prs):
    """Slide 6 – Liveness Detection"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_RED)
    add_text(slide, "👁️  ANTI-SPOOFING — LIVENESS DETECTION", 0.3, 0.2, 9, 0.75,
             font_size=26, bold=True, color=TEXT_WHITE)

    # EAR formula card
    add_rect(slide, 0.25, 1.28, 9.5, 1.3, BG_CARD, line_color=ACCENT_RED, line_width=Pt(1))
    add_text(slide, "Eye Aspect Ratio (EAR) Algorithm — Soukupová & Čech, 2016",
             0.45, 1.33, 9, 0.45, font_size=13, bold=True, color=ACCENT_RED)
    add_text(slide, "EAR  =  ( ‖p₂−p₆‖ + ‖p₃−p₅‖ )  /  ( 2 × ‖p₁−p₄‖ )",
             0.45, 1.75, 9, 0.6, font_size=18, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, font_name="Courier New")

    # How it works
    steps = [
        ("1️⃣", "Detect Face", "BlazeFace locates the face bounding box in the live video frame."),
        ("2️⃣", "Extract Landmarks", "FaceLandmarker returns 478 precise 3D landmarks — 12 used per eye."),
        ("3️⃣", "Compute EAR", "Vertical distances divided by horizontal eye width — drops sharply on blink."),
        ("4️⃣", "Count Blinks", "System requires 2 valid blinks within 8 seconds to pass the liveness gate."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        col = i % 2
        row_i = i // 2
        x = 0.25 + col * 4.9
        y = 2.78 + row_i * 1.55
        add_rect(slide, x, y, 4.65, 1.35, BG_CARD)
        add_rect(slide, x, y, 0.1, 1.35, ACCENT_RED)
        add_text(slide, num + "  " + title, x + 0.2, y + 0.08, 4.3, 0.48,
                 font_size=13, bold=True, color=TEXT_WHITE)
        add_text(slide, desc, x + 0.2, y + 0.55, 4.3, 0.75,
                 font_size=11, color=TEXT_LIGHT)

    # Config params
    add_rect(slide, 0.25, 5.95, 9.5, 0.72, RGBColor(0x2A, 0x0D, 0x1B))
    params = "EAR_THRESHOLD = 0.21     |     REQUIRED_BLINKS = 2     |     TIMEOUT = 8 seconds"
    add_text(slide, "⚙  " + params, 0.45, 6.02, 9.1, 0.55,
             font_size=12, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)


def slide7_gui(prs):
    """Slide 7 – GUI Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_TEAL)
    add_text(slide, "🎨  DESKTOP GUI — THREE-TAB INTERFACE", 0.3, 0.2, 9, 0.75,
             font_size=26, bold=True, color=TEXT_WHITE)

    tabs = [
        ("➕  Register", ACCENT_TEAL, [
            "Enter full name of the individual",
            "Start webcam — position face in frame",
            "Capture multiple samples (5-10 recommended)",
            "Click Register to persist to SQLite",
            "Multiple embeddings = higher accuracy",
        ]),
        ("👁  Recognize", ACCENT_GOLD, [
            "Click Start to begin live recognition feed",
            "BlazeFace detects faces in real time",
            "Liveness gate: 2 blinks within 8 seconds",
            "Identified persons shown with confidence %",
            "60-second cooldown per person per session",
        ]),
        ("🗂  Manage", ACCENT_GREEN, [
            "View all registered persons in table",
            "Embedding count per individual displayed",
            "Search and filter by name",
            "Delete individuals from the database",
            "Instant refresh after changes",
        ]),
    ]

    for i, (title, color, bullets) in enumerate(tabs):
        x = 0.25 + i * 3.25
        add_rect(slide, x, 1.28, 3.1, 5.35, BG_CARD, line_color=color, line_width=Pt(1.5))
        add_rect(slide, x, 1.28, 3.1, 0.55, color)
        add_text(slide, title, x + 0.1, 1.3, 2.9, 0.5,
                 font_size=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        for j, bullet in enumerate(bullets):
            add_text(slide, "•  " + bullet, x + 0.15, 1.98 + j * 0.72, 2.8, 0.65,
                     font_size=11, color=TEXT_LIGHT)

    # Bottom callout
    add_rect(slide, 0.25, 6.72, 9.5, 0.65, RGBColor(0x0A, 0x24, 0x38))
    add_text(slide, "🖥  Built with Tkinter 3 — Dark-themed, responsive, runs fully offline on Windows 10/11",
             0.45, 6.78, 9.1, 0.5, font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


def slide8_tech_stack(prs):
    """Slide 8 – Technology Stack & Dependencies"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_GOLD)
    add_text(slide, "📦  TECHNOLOGY STACK", 0.3, 0.2, 9, 0.75,
             font_size=28, bold=True, color=TEXT_WHITE)

    # Category cards
    categories = [
        ("🧠  Deep Learning", ACCENT_TEAL, [
            "PyTorch 2.11+  (inference backend)",
            "torchvision 0.26+ (transforms)",
            "facenet-pytorch 2.6 (InceptionResnetV1)",
        ]),
        ("👁  Vision & Detection", ACCENT_GOLD, [
            "MediaPipe 0.10+ (BlazeFace + Landmarks)",
            "OpenCV 4.13+ (camera capture & processing)",
            "478-point facial landmark model (TFLite)",
        ]),
        ("💾  Data & Storage", ACCENT_GREEN, [
            "SQLite 3 (zero-config embedded DB)",
            "NumPy 2.4+ (vector math)",
            "openpyxl 3.1+ (XLSX report generation)",
        ]),
        ("🖥  Interface & Utilities", ACCENT_RED, [
            "Tkinter 3 (cross-platform GUI)",
            "Pillow 12+ (image display in GUI)",
            "tqdm + requests (progress bars, HTTP)",
        ]),
    ]

    for i, (cat, color, items) in enumerate(categories):
        col = i % 2
        row_i = i // 2
        x = 0.25 + col * 4.88
        y = 1.28 + row_i * 2.42
        add_rect(slide, x, y, 4.65, 2.25, BG_CARD, line_color=color, line_width=Pt(1))
        add_rect(slide, x, y, 4.65, 0.52, color)
        add_text(slide, cat, x + 0.12, y + 0.06, 4.4, 0.45,
                 font_size=13, bold=True, color=BG_DARK)
        for j, item in enumerate(items):
            add_text(slide, "▸  " + item, x + 0.15, y + 0.62 + j * 0.52, 4.35, 0.46,
                     font_size=11, color=TEXT_LIGHT)

    # Python version badge
    add_rect(slide, 0.25, 6.18, 9.5, 0.55, RGBColor(0x14, 0x3A, 0x1E))
    add_text(slide, "🐍  Developed on Python 3.14  |  Tested on Windows 10 / 11  |  CPU & GPU compatible",
             0.45, 6.24, 9.1, 0.42, font_size=12, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


def slide9_results(prs):
    """Slide 9 – Results & Features Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 10, 1.1, BG_CARD)
    accent_bar(slide, 0, 1.1, 10, 0.06, ACCENT_GREEN)
    add_text(slide, "📊  RESULTS & KEY ACHIEVEMENTS", 0.3, 0.2, 9, 0.75,
             font_size=28, bold=True, color=TEXT_WHITE)

    # Metric tiles
    metrics = [
        ("512-dim", "FaceNet Vector", ACCENT_TEAL),
        ("478 pts", "Face Landmarks", ACCENT_GOLD),
        ("~30 ms", "Detection Latency", ACCENT_GREEN),
        ("2 Blinks", "Liveness Threshold", ACCENT_RED),
        ("8 sec", "Liveness Window", ACCENT_GOLD),
        ("60 sec", "Attendance Cooldown", ACCENT_TEAL),
    ]

    for i, (val, label, color) in enumerate(metrics):
        col = i % 3
        row_i = i // 3
        x = 0.25 + col * 3.22
        y = 1.28 + row_i * 1.48
        add_rect(slide, x, y, 3.05, 1.28, BG_CARD)
        add_rect(slide, x, y + 1.1, 3.05, 0.18, color)
        add_text(slide, val, x, y + 0.1, 3.05, 0.7,
                 font_size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, y + 0.78, 3.05, 0.35,
                 font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Feature checklist
    add_rect(slide, 0.25, 4.4, 9.5, 2.25, BG_CARD)
    add_text(slide, "✨  Feature Highlights", 0.45, 4.45, 4, 0.45,
             font_size=14, bold=True, color=ACCENT_TEAL)

    features_l = [
        "✔  Fully offline — zero cloud dependency",
        "✔  Auto model download on first run",
        "✔  Daily CSV + XLSX attendance reports",
        "✔  Asynchronous log writes (no UI freeze)",
    ]
    features_r = [
        "✔  Multi-capture registration for accuracy",
        "✔  Cosine similarity with confidence score",
        "✔  SQLite BLOB storage for embeddings",
        "✔  Headless test suite included",
    ]

    add_multiline_text(slide, features_l, 0.4, 4.92, 4.6, 1.7,
                       font_size=12, color=ACCENT_GREEN)
    add_multiline_text(slide, features_r, 5.1, 4.92, 4.6, 1.7,
                       font_size=12, color=ACCENT_GREEN)


def slide10_conclusion(prs):
    """Slide 10 – Conclusion & Future Work"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Full-width gradient header
    gradient_rect(slide, 0, 0, 10, 1.85, RGBColor(0x00, 0x44, 0x3C), BG_DARK)
    accent_bar(slide, 0, 1.85, 10, 0.06, ACCENT_TEAL)
    add_text(slide, "🏁  CONCLUSION & FUTURE SCOPE", 0.3, 0.2, 9, 0.75,
             font_size=28, bold=True, color=TEXT_WHITE)
    add_text(slide, "A complete, production-ready offline attendance system built with state-of-the-art CV/DL.",
             0.3, 0.9, 9.4, 0.55, font_size=13, italic=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Conclusion points
    add_rect(slide, 0.25, 2.05, 4.62, 3.1, BG_CARD)
    add_rect(slide, 0.25, 2.05, 4.62, 0.48, ACCENT_TEAL)
    add_text(slide, "🎯  Conclusion", 0.4, 2.08, 4.3, 0.42,
             font_size=14, bold=True, color=BG_DARK)
    conclusion = [
        "Successfully eliminates proxy attendance fraud",
        "Anti-spoofing prevents photo/screen attacks",
        "Full data privacy — 100% local, no cloud",
        "Clean 3-tab GUI for non-technical users",
        "Modular codebase — easy to extend",
        "Automated XLSX/CSV report generation",
    ]
    for i, pt in enumerate(conclusion):
        add_text(slide, "▸  " + pt, 0.4, 2.62 + i * 0.42, 4.35, 0.4,
                 font_size=11, color=TEXT_LIGHT)

    # Future work
    add_rect(slide, 5.13, 2.05, 4.62, 3.1, BG_CARD)
    add_rect(slide, 5.13, 2.05, 4.62, 0.48, ACCENT_GOLD)
    add_text(slide, "🚀  Future Scope", 5.28, 2.08, 4.3, 0.42,
             font_size=14, bold=True, color=BG_DARK)
    future = [
        "GPU acceleration for large-scale deployments",
        "Web dashboard for remote administration",
        "Mobile app integration (Android/iOS)",
        "Multi-camera simultaneous tracking",
        "Emotion / fatigue detection module",
        "Federated learning for privacy-preserving updates",
    ]
    for i, pt in enumerate(future):
        add_text(slide, "▸  " + pt, 5.28, 2.62 + i * 0.42, 4.35, 0.4,
                 font_size=11, color=TEXT_LIGHT)

    # Thank you banner
    add_rect(slide, 0, 5.32, 10, 1.68, RGBColor(0x07, 0x12, 0x1E))
    gradient_rect(slide, 0, 5.32, 10, 0.06, ACCENT_TEAL, ACCENT_GOLD)

    add_text(slide, "⬡  JAGTE RAHO", 0, 5.5, 10, 0.62,
             font_size=30, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "\"Stay Vigilant\"  —  Real-Time Face Recognition Attendance System",
             0, 6.12, 10, 0.45, font_size=13, italic=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "Thank you for your attention  ·  Questions welcome",
             0, 6.6, 10, 0.35, font_size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    slide1_title(prs)         ; print("  [1/10] Title slide")
    slide2_problem(prs)       ; print("  [2/10] Problem Statement")
    slide3_objectives(prs)    ; print("  [3/10] Objectives")
    slide4_architecture(prs)  ; print("  [4/10] System Architecture")
    slide5_face_recognition(prs) ; print("  [5/10] Face Recognition Pipeline")
    slide6_liveness(prs)      ; print("  [6/10] Liveness Detection")
    slide7_gui(prs)           ; print("  [7/10] GUI Overview")
    slide8_tech_stack(prs)    ; print("  [8/10] Technology Stack")
    slide9_results(prs)       ; print("  [9/10] Results & Achievements")
    slide10_conclusion(prs)   ; print("  [10/10] Conclusion & Future Work")

    out_path = "Jagte_Raho_Presentation.pptx"
    prs.save(out_path)
    print(f"\n✅  Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build_presentation()
